import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import csv
import logging

# Initialize logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class PowerPointProcess:

    @staticmethod
    def copy_and_rename_ppt(template_path, destination_path, entity_name, ppt_name):
        try:
            if not os.path.exists(template_path):
                raise FileNotFoundError(f"PowerPoint template not found at {template_path}")

            new_file_name = f"{entity_name}_{ppt_name}.pptx"
            new_file_path = os.path.join(destination_path, new_file_name)
            shutil.copy(template_path, new_file_path)

            logging.info(f"Template successfully copied and renamed to {new_file_path}")
            return new_file_path

        except FileNotFoundError as e:
            logging.error(f"Error: {e}")
        except Exception as e:
            logging.error(f"An unexpected error occurred: {e}")

    @staticmethod
    def replace_text_in_shape(shape, old_text, new_text):
        if not shape.has_text_frame:
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

    @staticmethod
    def replace_text_in_ppt(ppt_path, replacements):
        try:
            prs = Presentation(ppt_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    for old_text, new_text in replacements.items():
                        PowerPointProcess.replace_text_in_shape(shape, old_text, new_text)
            prs.save(ppt_path)
            logging.info(f"Text replaced in presentation: {ppt_path}")

        except Exception as e:
            logging.error(f"An error occurred: {e}")

    @staticmethod
    def add_and_format_table_in_powerpoint(ppt_path, slide_number, csv_path, left, top):
        prs = Presentation(ppt_path)

        with open(csv_path, newline='') as csvfile:
            reader = csv.reader(csvfile)
            data = list(reader)

        rows, cols = len(data), len(data[0])
        slide = prs.slides[slide_number - 1]

        # Dynamically adjust font size based on the number of rows
        font_size = Pt(9) if rows <= 10 else Pt(8) if rows <= 15 else Pt(7)

        # Calculate maximum width required for each column
        first_col_multiplier = 0.25
        default_col_multiplier = 0.0  # Adjusted multiplier
        col_max_widths = [
            Inches(first_col_multiplier * max(len(row[0]) for row in data) if i == 0 else
                default_col_multiplier * max(len(row[i]) for row in data))
            for i in range(cols)
        ]
        total_table_width = sum(col_max_widths)

        table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), total_table_width, Inches(3)).table

        for i, row in enumerate(data):
            for j, cell_text in enumerate(row):
                cell = table.cell(i, j)
                cell.width = col_max_widths[j]

                # Set cell to have no fill (transparent)
                cell.fill.background()

                para = cell.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = cell_text
                run.font.size = font_size

                # Align the first column left, and others right
                para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT

                # Make headers and "Total" cell bold and black
                if i == 0 or (i == rows - 1 and j == 0):
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Black color

        prs.save(ppt_path)
        logging.info(f"Table added to slide {slide_number} in PowerPoint file {ppt_path}")


    @staticmethod
    def add_tables_to_powerpoint(ppt_path, table_details):
        for slide_number, csv_path, left, top in table_details:
            try:
                if not os.path.exists(csv_path):
                    raise FileNotFoundError(f"CSV file not found: {csv_path}")
                PowerPointProcess.add_and_format_table_in_powerpoint(ppt_path, slide_number, csv_path, left, top)
                logging.info(f"Table added to slide {slide_number} from {csv_path}")

            except FileNotFoundError as e:
                logging.warning(f"Skipping table for slide {slide_number}: {e}")
            except Exception as e:
                logging.error(f"An error occurred while adding table to slide {slide_number}: {e}")

